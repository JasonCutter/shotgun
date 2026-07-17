from __future__ import annotations

import base64
import csv
import io
import json
import sys
from typing import Any


def selector_css(value: str) -> dict[str, Any]:
    return {"type": "CssSelector", "value": value}


def block(text: object, selectors: list[dict[str, Any]]) -> dict[str, Any] | None:
    value = " ".join(str(text).split())
    return {"text": value, "selectors": selectors} if value else None


def html_blocks(data: bytes) -> list[dict[str, Any]]:
    from bs4 import BeautifulSoup, Comment

    soup = BeautifulSoup(data.decode("utf-8"), "html.parser")
    for tag in soup.find_all(["script", "style", "noscript", "svg", "iframe", "canvas", "video", "audio", "object", "embed"]):
        tag.decompose()
    for comment in soup.find_all(string=lambda value: isinstance(value, Comment)):
        comment.extract()
    output: list[dict[str, Any]] = []
    tracked = soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "blockquote", "pre"])

    def css_path(tag: Any) -> str:
        parts: list[str] = []
        current = tag
        while getattr(current, "name", None) and current.name != "[document]":
            position = len(current.find_previous_siblings(current.name)) + 1
            parts.append(f"{current.name}:nth-of-type({position})")
            current = current.parent
        return " > ".join(reversed(parts))

    for tag in tracked:
        value = " ".join(tag.get_text(" ", strip=True).split())
        selector = css_path(tag)
        if soup.select_one(selector) is not tag:
            raise ValueError("HTML CSS selector did not round-trip")
        item = block(value, [selector_css(selector)])
        if item:
            output.append(item)
    if not output:
        item = block(" ".join(soup.get_text(" ", strip=True).split()), [selector_css("body")])
        if item:
            output.append(item)
    return output


def pdf_blocks(data: bytes) -> list[dict[str, Any]]:
    import pdfplumber

    if b"/Encrypt" in data:
        raise PermissionError("encrypted PDF")
    output: list[dict[str, Any]] = []
    with pdfplumber.open(io.BytesIO(data)) as document:
        if document.metadata.get("Encrypted") is True:
            raise PermissionError("encrypted PDF")
        for page_number, page in enumerate(document.pages, 1):
            for word in page.extract_words(use_text_flow=True, keep_blank_chars=False):
                item = block(
                    word["text"],
                    [
                        {"type": "PageSelector", "page": page_number},
                        {
                            "type": "BoundingBoxSelector",
                            "page": page_number,
                            "x": float(word["x0"]),
                            "y": float(word["top"]),
                            "width": float(word["x1"] - word["x0"]),
                            "height": float(word["bottom"] - word["top"]),
                            "unit": "pt",
                        },
                    ],
                )
                if item:
                    output.append(item)
    return output


def docx_blocks(data: bytes) -> list[dict[str, Any]]:
    from docx import Document

    document = Document(io.BytesIO(data))
    output: list[dict[str, Any]] = []
    for index, paragraph in enumerate(document.paragraphs, 1):
        item = block(paragraph.text, [selector_css(f"word/paragraph[{index}]")])
        if item:
            output.append(item)
    for table_index, table in enumerate(document.tables, 1):
        for row_index, row in enumerate(table.rows, 1):
            for column_index, cell in enumerate(row.cells, 1):
                item = block(
                    cell.text,
                    [
                        {
                            "type": "CellSelector",
                            "sheet": f"table-{table_index}",
                            "cell": f"R{row_index}C{column_index}",
                            "row": row_index,
                            "column": column_index,
                        }
                    ],
                )
                if item:
                    output.append(item)
    return output


def xlsx_blocks(data: bytes) -> list[dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=False)
    output: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell_value in row:
                if cell_value.value is None:
                    continue
                item = block(
                    cell_value.value,
                    [
                        {
                            "type": "CellSelector",
                            "sheet": sheet.title,
                            "cell": cell_value.coordinate,
                            "row": cell_value.row,
                            "column": cell_value.column,
                        }
                    ],
                )
                if item:
                    output.append(item)
    workbook.close()
    return output


def csv_blocks(data: bytes) -> list[dict[str, Any]]:
    from openpyxl.utils import get_column_letter

    output: list[dict[str, Any]] = []
    for row_number, row in enumerate(csv.reader(io.StringIO(data.decode("utf-8"))), 1):
        for column_number, value in enumerate(row, 1):
            item = block(
                value,
                [
                    {
                        "type": "CellSelector",
                        "sheet": "CSV",
                        "cell": f"{get_column_letter(column_number)}{row_number}",
                        "row": row_number,
                        "column": column_number,
                    }
                ],
            )
            if item:
                output.append(item)
    return output


def pptx_blocks(data: bytes) -> list[dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    output: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            text_value = getattr(shape, "text", "")
            item = block(
                text_value,
                [
                    {"type": "ShapeSelector", "slide": slide_number, "shapeId": str(shape.shape_id)},
                    {
                        "type": "BoundingBoxSelector",
                        "page": slide_number,
                        "x": float(shape.left / 12700),
                        "y": float(shape.top / 12700),
                        "width": float(shape.width / 12700),
                        "height": float(shape.height / 12700),
                        "unit": "pt",
                    },
                ],
            )
            if item:
                output.append(item)
    return output


def image_blocks(data: bytes, description: str | None) -> list[dict[str, Any]]:
    from PIL import Image

    with Image.open(io.BytesIO(data)) as image:
        image.verify()
    with Image.open(io.BytesIO(data)) as image:
        width, height = image.size
    if not description:
        raise RuntimeError("MULTIMODAL_VALIDATION_REQUIRED")
    item = block(
        description,
        [{"type": "BoundingBoxSelector", "x": 0, "y": 0, "width": width, "height": height, "unit": "px"}],
    )
    return [item] if item else []


def main() -> None:
    request = json.load(sys.stdin)
    data = base64.b64decode(request["contentBase64"], validate=True)
    media_type = request["mediaType"]
    handlers = {
        "text/html": lambda: html_blocks(data),
        "application/pdf": lambda: pdf_blocks(data),
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": lambda: docx_blocks(data),
        "text/csv": lambda: csv_blocks(data),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": lambda: xlsx_blocks(data),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": lambda: pptx_blocks(data),
        "image/png": lambda: image_blocks(data, request.get("imageDescription")),
        "image/jpeg": lambda: image_blocks(data, request.get("imageDescription")),
    }
    if media_type not in handlers:
        raise NotImplementedError(media_type)
    blocks = handlers[media_type]()
    if not blocks:
        raise ValueError("document contains no accessible text")
    json.dump({"status": "OK", "blocks": blocks}, sys.stdout, ensure_ascii=False)


try:
    main()
except NotImplementedError as error:
    json.dump({"status": "FORMAT_UNSUPPORTED", "message": str(error)}, sys.stdout)
except PermissionError as error:
    json.dump({"status": "FORMAT_ENCRYPTED", "message": str(error)}, sys.stdout)
except RuntimeError as error:
    status = str(error) if str(error) == "MULTIMODAL_VALIDATION_REQUIRED" else "FORMAT_CORRUPT"
    json.dump({"status": status, "message": str(error)}, sys.stdout)
except Exception as error:
    message = str(error)
    lowered = f"{error.__class__.__name__} {message}".lower()
    status = "FORMAT_ENCRYPTED" if "password" in lowered or "encrypted" in lowered else "FORMAT_CORRUPT"
    json.dump({"status": status, "message": message or error.__class__.__name__}, sys.stdout)
